import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_pptx_from_md(md_file_path, pptx_file_path):
    print(f"Converting {md_file_path} to {pptx_file_path}...")
    
    prs = Presentation()
    
    with open(md_file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split content by slide headers (e.g., ## 幻灯片1：封面)
    slides_content = re.split(r'## 幻灯片\d+：', content)
    
    # The first part is usually the main title (e.g., # 第1章...)
    if slides_content:
        main_title_match = re.search(r'# (.*?)\n', slides_content[0])
        main_title = main_title_match.group(1) if main_title_match else "Python Data Analysis"
        slides_content = slides_content[1:] # Skip the header part
    
    for slide_text in slides_content:
        lines = slide_text.strip().split('\n')
        if not lines:
            continue
            
        title_line = lines[0].strip()
        body_lines = lines[1:]
        
        # Add slide
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title_line
        
        # Set body
        if body_lines:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.word_wrap = True
            
            p = None
            is_code_block = False
            
            for line in body_lines:
                line = line.strip()
                if not line:
                    continue
                
                # Handle code blocks
                if line.startswith('```'):
                    is_code_block = not is_code_block
                    continue
                
                if p is None:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # Check for bullet points
                if line.startswith('- '):
                    p.text = line[2:]
                    p.level = 0
                elif line.startswith('  - '):
                    p.text = line[4:]
                    p.level = 1
                else:
                    p.text = line
                
                # Formatting for code
                if is_code_block:
                    p.font.name = 'Courier New'
                    p.font.size = Pt(14)
                else:
                    p.font.size = Pt(18)

    prs.save(pptx_file_path)
    print(f"Saved: {pptx_file_path}")

def main():
    directory = "/Users/zhaoqinhuang/david_project/course/ppt_outlines"
    files = [f for f in os.listdir(directory) if (f.startswith("chapter") or f.startswith("course_ppt")) and f.endswith(".md")]
    
    for file in files:
        md_path = os.path.join(directory, file)
        pptx_name = file.replace(".md", ".pptx")
        pptx_path = os.path.join(directory, pptx_name)
        create_pptx_from_md(md_path, pptx_path)

if __name__ == "__main__":
    main()
